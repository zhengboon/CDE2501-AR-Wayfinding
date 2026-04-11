"""
CDE2501 Final Presentation — Solution: AR Wayfinding App
Follows rubric: Site Analysis (20%), Problem ID & Analysis (25%),
Framing (25%), Solution Evaluation (20%), Presentation (10%)
"""
import subprocess, sys
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_COLOR = RGBColor(0x0E, 0xA5, 0xE9)
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def add_slide(title_text, bullets=None, subtitle=None, layout_idx=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    title = slide.shapes.title
    title.text = title_text
    for para in title.text_frame.paragraphs:
        para.font.size = Pt(36)
        para.font.bold = True
        para.font.color.rgb = TITLE_COLOR

    if subtitle:
        body = slide.placeholders[1]
        body.text = subtitle
        for para in body.text_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = BODY_COLOR

    if bullets:
        body = slide.placeholders[1]
        body.text = ""
        for i, bullet in enumerate(bullets):
            if i == 0:
                body.text = bullet
                body.text_frame.paragraphs[0].font.size = Pt(20)
                body.text_frame.paragraphs[0].font.color.rgb = BODY_COLOR
            else:
                p = body.text_frame.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)
                p.font.color.rgb = BODY_COLOR
                p.space_before = Pt(8)
    return slide

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "AR Wayfinding for Queenstown's Elderly"
subtitle.text = "CDE2501 Sustainable Systems for Liveable Cities\nGroup 2 — Queenstown\nAY 2025-2026 Semester 2"
for para in title.text_frame.paragraphs:
    para.font.size = Pt(40)
    para.font.bold = True

# ============================================================
# SLIDE 2: Site Analysis — Context
# ============================================================
add_slide("Site Analysis: Queenstown", [
    "Singapore's first satellite town — developed in the 1950s",
    "Mature estate with over 70 years of history",
    "Super-aged: 22.6% of residents are 65+ years old",
    "Part of HDB's Health District @ Queenstown initiative",
    "NUS-NUHS collaboration area for healthy ageing research",
    "URA Master Plan 2019: increased redevelopment, stronger MRT planning, healthcare emphasis",
])

# ============================================================
# SLIDE 3: Site Analysis — Physical Characteristics
# ============================================================
add_slide("Site Analysis: Physical Environment", [
    "Site area: 159,301 m² | Perimeter: 1,645 m",
    "Elevation range: Min 6.6m, Median 15.3m, Max 24.76m (8.7m difference)",
    "Hilly terrain with constant uphill routes between blocks",
    "Older flat layouts with narrow corridors (insufficient for wheelchairs)",
    "Limited or poorly sheltered walkways between blocks",
    "Cluttered common areas — overflowing recycling, motorcycles on walkways",
    "Dementia Wayfinding Murals at Stirling View (existing community effort)",
])

# ============================================================
# SLIDE 4: Site Analysis — Demographics & Stakeholders
# ============================================================
add_slide("Site Analysis: Who Lives Here", [
    "22.6% aged 65+ — classified as super-aged estate",
    "65.6% working adults (15-64 years)",
    "Key landmarks: Queenstown MRT, Mei Ling Market & Food Centre, Queenstown Polyclinic",
    "Lions Befrienders and Care Corner — active social service presence",
    "Mixed land use: residential, commercial, health & medical, religious, recreational",
    "Strong MRT connectivity but poor last-mile within the estate",
])

# ============================================================
# SLIDE 5: Problem Statement
# ============================================================
add_slide("Problem Statement", [
    '"In Queenstown\'s interconnected HDB blocks, inconsistent level connections',
    'and poor spatial legibility produce unsafe and disorienting navigation',
    'conditions that disproportionately impact wheelchair users, dementia patients,',
    'cane users, and medically challenged elderly residents, heightening fall risk',
    'and undermining barrier-free mobility."',
    "",
    "Root causes: 8.7m elevation change, limited ramps, poor signage, narrow corridors",
])

# ============================================================
# SLIDE 6: Problem Analysis — Affected Communities
# ============================================================
add_slide("Problem Analysis: Who Is Affected", [
    "Wheelchair users: 8.7m elevation climb, no nearby ramps, constant uphill → limited independence",
    "Dementia patients: repetitive layout, long detour routes → increased confusion, anxiety",
    "Cane/walking aid users: steep terrain, stair-heavy routes → higher fall risk, fatigue",
    "Medically challenged elderly: exhausting travel routes → avoid leaving home, social isolation",
    "Caregivers: physical strain assisting uphill, emotional stress burden",
    "Town Council: resident complaints, upgrade cost pressure, safety liability",
])

# ============================================================
# SLIDE 7: Systems Thinking — Causal Loop
# ============================================================
add_slide("Systems Thinking: Causal Loop Diagram", [
    "Reinforcing Loop (R1): Mobility Barriers → Risk of Safety Hazards → Reduced Use of Public Spaces",
    "  → Less Local Stewardship → More Clutter → More Passive Common Areas → More Barriers",
    "",
    "Balancing Loop (B1): Mobility Barriers → Demand for Wayfinding Improvements",
    "  → Better Signage → Reduced Safety Hazards → Improved Quality of Life",
    "",
    "Our intervention targets B1: improving wayfinding effectiveness to reduce safety hazards",
])

# ============================================================
# SLIDE 8: Framing — Why Wayfinding
# ============================================================
add_slide("Framing: Why Digital Wayfinding", [
    "Physical infrastructure (ramps, lifts) costs $200-400K per installation, takes years",
    "Signage helps but is static — doesn't adapt to individual needs or conditions",
    "22.6% elderly population needs guidance NOW, not in 5 years",
    "Digital wayfinding: zero infrastructure cost, instant deployment, personalized routes",
    "Collects walking data that INFORMS future physical infrastructure investment",
    "Aligns with Health District @ Queenstown vision of technology-enabled healthy ageing",
])

# ============================================================
# SLIDE 9: Solution — The AR Wayfinding App
# ============================================================
add_slide("Solution: AR Wayfinding App", [
    "Flight-tracker AR view: point phone at destination → see name, distance, ETA on camera",
    "Safety-weighted routing: avoids stairs for wheelchairs, steep slopes for elderly",
    "Minimap with heading arrow: always know which direction you're facing",
    "Two profiles: Elderly (1.0 m/s) and Wheelchair (0.8 m/s) with rain mode",
    "Works on any Android phone with GPS + compass — no special hardware needed",
    "Data syncs from Google Drive — updates without reinstalling the app",
])

# ============================================================
# SLIDE 10: Solution — Data Collection
# ============================================================
add_slide("Solution: Crowd-Sourced Walking Data", [
    "Record labeled paths: testers walk from A to B, GPS breadcrumbs saved every 1 second",
    "Auto-screenshots on heading changes (>30°) + every 10 seconds + manual Snap",
    "Telemetry CSV: GPS, heading, altitude, floor estimate, GPS accuracy, loss detection",
    "Saved paths show where elderly ACTUALLY walk vs. where planners think they walk",
    "Data feeds back to HDB: evidence for ramp placement, walkway sheltering, signage",
    "Share button sends telemetry via Drive/email — no technical knowledge needed",
])

# ============================================================
# SLIDE 11: Solution — How It Works
# ============================================================
add_slide("How It Works: User Journey", [
    "1. Open app → data downloads from Google Drive (3MB, takes seconds)",
    "2. Select destination: 'Mei Ling Hawker Centre'",
    "3. App calculates safest route avoiding stairs and steep slopes",
    "4. AR view shows direction arrow on camera — follow it",
    "5. Minimap shows progress with heading indicator",
    "6. Tap 'Record Path' to save the walk for future graph improvement",
    "7. Arrive. Total guided time shown. Path saved for analysis.",
])

# ============================================================
# SLIDE 12: Solution Evaluation — Stakeholders
# ============================================================
add_slide("Solution Evaluation: Stakeholder Impact", [
    "Wheelchair users: routes avoid stairs entirely → independent mobility restored",
    "Dementia patients: AR arrows + familiar landmarks → reduced confusion",
    "Cane users: flat/gentle routes prioritized → lower fall risk",
    "Caregivers: can send elderly independently → reduced physical & emotional burden",
    "HDB/Town Council: real walking data → evidence-based infrastructure planning",
    "Health District @ Queenstown: technology-enabled ageing in place",
])

# ============================================================
# SLIDE 13: Solution Evaluation — Trade-offs
# ============================================================
add_slide("Solution Evaluation: Trade-offs & Limitations", [
    "Requires smartphone: many elderly (70-90 years) lack digital literacy",
    "  → Mitigation: family/caregiver sets up the app once, elderly follows AR arrows",
    "GPS unreliable between tall blocks: 10-30m drift in signal shadow zones",
    "  → Mitigation: dead reckoning fallback planned (compass + step counting)",
    "Can't fix physical terrain: if no ramp exists, no alternative path",
    "  → Mitigation: walking data identifies exactly WHERE ramps are needed most",
    "No indoor navigation yet: corridors and lift lobbies need BLE/WiFi positioning",
])

# ============================================================
# SLIDE 14: Social, Environmental, Economic Impact
# ============================================================
add_slide("Impact: Social, Environmental, Economic", [
    "Social: reduces isolation (elderly leave home more), builds independence, dementia-friendly",
    "Environmental: encourages walking over motorized transport, zero emissions from app",
    "Economic: $0 infrastructure cost vs. $200-400K per ramp installation",
    "  App deployment: free (Google Drive hosting, no server needed)",
    "  Data collection: replaces expensive walking audits ($50-100K consultant studies)",
    "  Scalable: same app works in any HDB estate — just update the graph data",
])

# ============================================================
# SLIDE 15: Conclusion
# ============================================================
add_slide("Conclusion", [
    "Queenstown's elderly face unsafe, exhausting, disorienting daily navigation",
    "Physical fixes (ramps, lifts) are expensive and slow — 5+ year timeline",
    "AR wayfinding provides immediate, personalized, data-collecting guidance",
    "Every walk recorded improves the system — crowd-sourced path optimization",
    "The app doesn't replace infrastructure — it makes infrastructure investment smarter",
    "",
    "Digital wayfinding companion for the Health District @ Queenstown initiative",
])

output = "CDE2501_Final_Presentation_App_Solution.pptx"
prs.save(output)
print(f"Saved: {output}")
