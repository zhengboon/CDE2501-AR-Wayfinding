"""
CDE2501 Final Presentation — Solution: AR Wayfinding App + Shuttle Bus
Follows rubric: Site Analysis (20%), Problem ID & Analysis (25%),
Framing (25%), Solution Evaluation (20%), Presentation (10%)
"""
import subprocess, sys
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)

def add_slide(title_text, bullets=None, layout_idx=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    title = slide.shapes.title
    title.text = title_text
    for para in title.text_frame.paragraphs:
        para.font.size = Pt(36)
        para.font.bold = True
        para.font.color.rgb = TITLE_COLOR
    if bullets:
        body = slide.placeholders[1]
        body.text = ""
        for i, bullet in enumerate(bullets):
            if i == 0:
                body.text = bullet
                body.text_frame.paragraphs[0].font.size = Pt(20)
                body.text_frame.paragraphs[0].font.color.rgb = BODY_COLOR
            else:
                p = body.text_frame.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)
                p.font.color.rgb = BODY_COLOR
                p.space_before = Pt(8)
    return slide

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "AR Wayfinding + Intra-Estate Shuttle\nfor Queenstown's Elderly"
subtitle.text = "CDE2501 Sustainable Systems for Liveable Cities\nGroup 2 — Queenstown\nAY 2025-2026 Semester 2"
for para in title.text_frame.paragraphs:
    para.font.size = Pt(36)
    para.font.bold = True

# ============================================================
# SLIDE 2: Site Analysis — Context
# ============================================================
add_slide("Site Analysis: Queenstown", [
    "Singapore's first satellite town — developed in the 1950s",
    "Mature estate with over 70 years of history",
    "Super-aged: 22.6% of residents are 65+ years old",
    "Part of HDB's Health District @ Queenstown initiative",
    "NUS-NUHS collaboration area for healthy ageing research",
    "URA Master Plan 2019: increased redevelopment, stronger MRT planning, healthcare emphasis",
])

# ============================================================
# SLIDE 3: Site Analysis — Physical Environment
# ============================================================
add_slide("Site Analysis: Physical Environment", [
    "Site area: 159,301 m² | Perimeter: 1,645 m",
    "Elevation range: Min 6.6m, Median 15.3m, Max 24.76m (8.7m difference)",
    "Hilly terrain with constant uphill routes between blocks",
    "Older flat layouts with narrow corridors (insufficient for wheelchairs)",
    "Limited or poorly sheltered walkways between blocks",
    "9-minute walk from hilltop blocks to MRT — exhausting for elderly",
    "Lack of ramps or lifts makes it incredibly difficult for limited mobility",
])

# ============================================================
# SLIDE 4: Site Analysis — Demographics
# ============================================================
add_slide("Site Analysis: Who Lives Here", [
    "22.6% aged 65+ — classified as super-aged estate",
    "65.6% working adults (15-64 years)",
    "Key landmarks: Queenstown MRT, Mei Ling Market & Food Centre, Queenstown Polyclinic",
    "Lions Befrienders and Care Corner — active social service presence",
    "Strong MRT connectivity but poor last-mile within the estate",
    "Medically challenged elderly avoid leaving home due to exhausting routes",
])

# ============================================================
# SLIDE 5: Problem Statement
# ============================================================
add_slide("Problem Statement", [
    '"In Queenstown\'s interconnected HDB blocks, inconsistent level connections',
    'and poor spatial legibility produce unsafe and disorienting navigation',
    'conditions that disproportionately impact wheelchair users, dementia patients,',
    'cane users, and medically challenged elderly residents, heightening fall risk',
    'and undermining barrier-free mobility."',
    "",
    "Two root causes: PHYSICAL barrier (8.7m elevation, no ramps) + COGNITIVE barrier (poor signage, confusing layout)",
])

# ============================================================
# SLIDE 6: Problem Analysis — Affected Communities
# ============================================================
add_slide("Problem Analysis: Who Is Affected", [
    "Wheelchair users: 8.7m elevation climb, no nearby ramps → limited independence, greater caregiver reliance",
    "Dementia patients: repetitive layout, long detours → increased confusion, higher anxiety",
    "Cane/walking aid users: steep terrain, stair-heavy routes → higher fall risk, avoid essential trips",
    "Medically challenged elderly: exhausting routes → avoid leaving home, reduced healthcare access, social isolation",
    "Caregivers: physical strain pushing uphill, emotional stress burden",
    "Town Council: resident complaints, upgrade cost pressure, safety liability risks",
])

# ============================================================
# SLIDE 7: Systems Thinking — Causal Loop
# ============================================================
add_slide("Systems Thinking: Causal Loop Diagram", [
    "Reinforcing Loop (R1): Mobility Barriers → Safety Hazards → Reduced Use of Public Spaces",
    "  → Less Stewardship → More Clutter → More Passive Areas → Worse Barriers",
    "",
    "Balancing Loop (B1): Barriers → Demand for Improvements → Better Wayfinding → Reduced Hazards",
    "",
    "Our solution targets BOTH loops:",
    "  Digital wayfinding (B1) — improves navigation effectiveness immediately",
    "  Shuttle bus (R1) — breaks the physical barrier that keeps elderly home",
])

# ============================================================
# SLIDE 8: Framing — Two Problems, Two Solutions
# ============================================================
add_slide("Framing: Two Problems Need Two Solutions", [
    "COGNITIVE barrier: 'I don't know which way is safe'",
    "  → Digital wayfinding with AR guidance, safety-weighted routing",
    "",
    "PHYSICAL barrier: 'I physically cannot climb this hill'",
    "  → Intra-estate shuttle eliminates the 8.7m climb entirely",
    "",
    "Neither solution alone is sufficient:",
    "  App without bus: routes around stairs, but if no flat path exists, elderly still can't move",
    "  Bus without app: gets elderly to the area, but last-mile from stop to destination still confusing",
])

# ============================================================
# SLIDE 9: Solution Part 1 — AR Wayfinding App
# ============================================================
add_slide("Solution Part 1: AR Wayfinding App", [
    "Flight-tracker AR view: point phone at destination → see name, distance, ETA on camera",
    "Safety-weighted routing: avoids stairs for wheelchairs, steep slopes for elderly",
    "Minimap with heading arrow: always know which direction you're facing",
    "Two profiles: Elderly (1.0 m/s) and Wheelchair (0.8 m/s) with rain mode",
    "Record walking paths: GPS breadcrumbs every 1s + auto-screenshots at turns",
    "Google Drive sync: data updates without reinstalling, telemetry shared from device",
    "Works on any Android phone — no special hardware needed",
])

# ============================================================
# SLIDE 10: Solution Part 2 — Intra-Estate Shuttle
# ============================================================
add_slide("Solution Part 2: Intra-Estate Shuttle Bus", [
    "Electric minibus (8-12 seats, wheelchair ramp, low-floor boarding)",
    "Circular loop: Queenstown MRT → Mei Ling Hawker → Hilltop Blocks → Stirling Road → MRT",
    "~1.6km route, 5-6 minute loop, 10-minute frequency",
    "Free for residents aged 65+ (subsidized by Town Council / Health District funding)",
    "4 stops at key locations: MRT (7m), Hawker (15m), Hilltop (22m), Stirling (8m)",
    "Precedent: HDB trialled intra-town shuttles in Teck Ghee and Bukit Panjang",
    "Eliminates the 8.7m climb — the biggest physical barrier in the estate",
])

# ============================================================
# SLIDE 11: How They Work Together
# ============================================================
add_slide("The Integrated System: App + Bus", [
    "Current journey: Block 163 → Hawker Centre = 9 min uphill, exhausting, confusing",
    "",
    "With integrated system:",
    "1. Open app, tap 'Mei Ling Hawker Centre'",
    "2. App says: 'Bus Stop B — 1 min walk, flat, sheltered'",
    "3. AR arrow guides to bus stop (60m, no stairs)",
    "4. Bus arrives in 4 min, ride 3 min",
    "5. Get off, AR guides 20m to hawker entrance",
    "6. Total: 2 min walking + 3 min bus. Zero stairs. Zero confusion.",
])

# ============================================================
# SLIDE 12: Data-Driven Bus Planning
# ============================================================
add_slide("App Data Informs Bus Planning", [
    "GPS trails show where elderly actually walk → identifies high-traffic routes",
    "Heading changes reveal confusion points → needs signage or landmarks",
    "Path avoidance patterns show which stairs/slopes are unused → ramp priority",
    "Time-of-day patterns reveal peak elderly movement hours → bus schedule optimization",
    "Wheelchair vs elderly profile usage → which routes need wheelchair ramps vs handrails",
    "",
    "The app is a free, continuous, crowd-sourced mobility survey",
    "HDB currently pays $50-100K for consultant walking audits — app does it automatically",
])

# ============================================================
# SLIDE 13: Solution Evaluation — Stakeholder Impact
# ============================================================
add_slide("Solution Evaluation: Stakeholder Impact", [
    "Wheelchair users: bus eliminates climb + app finds flat path to stop → full independence",
    "Dementia patients: AR guidance reduces confusion + bus reduces distance walked",
    "Cane users: bus handles steep segments + app avoids remaining stairs → lower fall risk",
    "Medically challenged: bus stop within 1 min of any block → no longer avoid leaving home",
    "Caregivers: can send elderly independently → reduced physical & emotional burden",
    "HDB/Town Council: real walking data → evidence-based infrastructure planning",
])

# ============================================================
# SLIDE 14: Solution Evaluation — Trade-offs
# ============================================================
add_slide("Solution Evaluation: Trade-offs", [
    "App requires smartphone: many elderly lack digital literacy",
    "  → Mitigation: family sets up once, elderly follows simple AR arrows",
    "Bus costs $500-700K Year 1 + $100-150K/year operating",
    "  → Mitigation: Health District funding; app data proves demand to justify investment",
    "Bus has fixed route — can't serve every block",
    "  → Mitigation: app handles last-mile from any block to the nearest stop",
    "Low ridership risk (precedent from other estates)",
    "  → Mitigation: app telemetry measures actual usage — adjust or cancel if underperforming",
])

# ============================================================
# SLIDE 15: Social, Environmental, Economic Impact
# ============================================================
add_slide("Impact: Social, Environmental, Economic", [
    "Social: reduces isolation, builds independence, dementia-friendly, inclusive for all mobility levels",
    "Environmental: electric bus = zero emissions; app encourages walking for able-bodied segments",
    "Economic: app costs $0 to operate; bus justified by data (not guesswork)",
    "  Combined: cheaper than building ramps everywhere ($200-400K each)",
    "  App data replaces expensive consultant walking audits",
    "  Scalable: same app + bus model can be replicated in other super-aged estates",
])

# ============================================================
# SLIDE 16: Comparison — App Only vs App + Bus
# ============================================================
add_slide("Why Both, Not Just One", [
    "App only: routes elderly around stairs for 9 min uphill — still exhausting",
    "Bus only: gets to the area, but last-mile is still confusing, no data collected",
    "App + Bus: 2 min flat walk + 3 min ride — zero stairs, zero confusion, data flowing",
    "",
    "The app makes the bus SMARTER (data-driven route/stop planning)",
    "The bus makes the app USEFUL (reduces 9 min walk to 1 min walk to stop)",
    "Together: no elderly resident is more than 2 minutes from accessible transport",
])

# ============================================================
# SLIDE 17: Conclusion
# ============================================================
add_slide("Conclusion", [
    "Queenstown's elderly face two barriers: physical (terrain) and cognitive (navigation)",
    "The AR wayfinding app solves the cognitive barrier immediately, at zero cost",
    "The intra-estate shuttle solves the physical barrier — eliminates the 8.7m climb",
    "The app's walking data justifies the bus investment with evidence, not assumptions",
    "Together: a complete, integrated mobility system for Singapore's first Health District",
    "",
    "Digital guidance + accessible transport = independent, confident, connected elderly",
])

output = "CDE2501_Final_Presentation_App_Bus_Solution.pptx"
prs.save(output)
print(f"Saved: {output}")
