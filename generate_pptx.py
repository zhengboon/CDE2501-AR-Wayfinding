import sys
import subprocess

def install_pptx():
    try:
        import pptx
    except ImportError:
        print("python-pptx not found. Installing...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx
    return pptx

pptx = install_pptx()
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Title slide layout
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "CDE2501 AR Wayfinding"
subtitle.text = "How the Program Works\nAn AR navigation MVP prioritizing elderly and wheelchair-accessible routing."

# Bullet slide layout
bullet_slide_layout = prs.slide_layouts[1]

# Slide 2: Project Overview
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "1. Project Overview"
tf = body_shape.text_frame
tf.text = "What it is: An Augmented Reality wayfinding app built in Unity."
p = tf.add_paragraph()
p.text = "Target Areas: Queenstown estate and NUS Engineering campus."
p = tf.add_paragraph()
p.text = "Goal: Provide safe, elderly-first, and barrier-free routing."
p = tf.add_paragraph()
p.text = "Cross-Platform: Targets Android (ARCore) and iOS (ARKit)."

# Slide 3: Core Features
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "2. Core Features"
tf = body_shape.text_frame
tf.text = "AR Guidance: 3D arrows overlaid on the camera feed."
p = tf.add_paragraph()
p.text = "Minimap Overlay: Top-down view highlighting the path."
p = tf.add_paragraph()
p.text = "Multi-Area Support: Dynamic switching between local map graphs."
p = tf.add_paragraph()
p.text = "Safety-Weighted Routing: Profiles for Elderly (1.0m/s) and Wheelchair (0.8m/s)."
p = tf.add_paragraph()
p.text = "Telemetry Recording: Live GPS/Compass logging to CSV."

# Slide 4: How Routing Works
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "3. How Routing Works (The Brain)"
tf = body_shape.text_frame
tf.text = "Algorithm: Weighted A* (A-Star) pathfinding on a node-edge graph."
p = tf.add_paragraph()
p.text = "Data: The graph (estate_graph.json) holds walkable paths."
p = tf.add_paragraph()
p.text = "Profiles: Edges are evaluated based on user Safety Profiles."
p = tf.add_paragraph()
p.text = "Execution: System anchors the user to the closest start node and computes the path to the destination."

# Slide 5: Data Generation Pipeline
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "4. Data Generation Pipeline (Offline)"
tf = body_shape.text_frame
tf.text = "Map Tiles: Fetched via OneMap API, Google Maps, or OSM."
p = tf.add_paragraph()
p.text = "Graph Gen: Python scripts (generate_osm_graph_from_kml.py) parse KML boundaries into interconnected nodes/edges."
p = tf.add_paragraph()
p.text = "Street View: Can download Street View panoramas corresponding to walking routes."

# Slide 6: Sensors & Locomotion
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "5. Sensors & Locomotion (Runtime)"
tf = body_shape.text_frame
tf.text = "Device Sensors: Native GPS and Compass heavily smoothed for AR stability."
p = tf.add_paragraph()
p.text = "Simulation Mode: WASD walking and camera rotation overrides sensors for laptop testing."
p = tf.add_paragraph()
p.text = "Auto-Refresh: If the user strays > 2.0m off-path, the route calculates a new path automatically."

# Slide 7: Telemetry & Field Testing
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "6. Telemetry & Field Testing"
tf = body_shape.text_frame
tf.text = "Built-in Recorder: Toggle Rec: ON/OFF in the UI."
p = tf.add_paragraph()
p.text = "Data Logged: Time, Lat, Lon, Heading, StartNode, Destination, RouteDistance into local device CSV files."
p = tf.add_paragraph()
p.text = "Post-analysis: Used to validate safety of routes and refine pathfinding edge weights based on real-world walking data."

# Slide 8: Summary of the User Journey
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "7. Summary of the User Journey"
tf = body_shape.text_frame
tf.text = "1. Launch: Load map graph into memory."
p = tf.add_paragraph()
p.text = "2. Context: Lock onto GPS/Compass; determine start node."
p = tf.add_paragraph()
p.text = "3. Selection: User picks destination."
p = tf.add_paragraph()
p.text = "4. Computation: A* determines safest path."
p = tf.add_paragraph()
p.text = "5. Guidance: Render AR arrows and Minimap progress."
p = tf.add_paragraph()
p.text = "6. Arrival: Complete at destination map node."

output_file = "CDE2501_AR_Wayfinding_Presentation.pptx"
prs.save(output_file)
print(f"Presentation saved to {output_file}")
